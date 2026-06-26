"""
文件解析器模块 - 支持多种文件类型的文本提取

支持的文件类型：
- txt: 纯文本文件
- md: Markdown 文件
- pdf: PDF 文件（每页独立）
- docx: Word 文件
- xlsx: Excel 文件（每个 Sheet 独立）
- pptx: PPT 文件（每页独立）
"""

from __future__ import annotations

import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from loguru import logger

from config import config_data as config


@dataclass
class ParsedDocument:
    """解析后的文档结构"""

    content: str  # 文本内容
    doc_type: str  # 文档类型: txt, md, pdf, docx, xlsx, pptx
    source: str  # 文件名
    # 可选的元数据
    page_number: int | None = None  # PDF/PPT 页码
    sheet_name: str | None = None  # Excel Sheet 名称
    sheet_index: int | None = None  # Excel Sheet 序号
    extra_metadata: dict[str, Any] | None = None  # 其他元数据


def parse_file(file_bytes: bytes, filename: str) -> list[ParsedDocument]:
    """
    解析文件内容，返回文档列表

    Args:
        file_bytes: 文件二进制内容
        filename: 文件名（用于判断类型）

    Returns:
        解析后的文档列表（PDF/PPT/XLSX 可能返回多个文档）
    """
    ext = Path(filename).suffix.lower()

    if ext == ".txt":
        return _parse_txt(file_bytes, filename)
    elif ext == ".md":
        return _parse_markdown(file_bytes, filename)
    elif ext == ".pdf":
        return _parse_pdf(file_bytes, filename)
    elif ext == ".docx":
        return _parse_docx(file_bytes, filename)
    elif ext == ".xlsx":
        return _parse_xlsx(file_bytes, filename)
    elif ext == ".pptx":
        return _parse_pptx(file_bytes, filename)
    else:
        raise ValueError(f"不支持的文件类型: {ext}")


def _parse_txt(file_bytes: bytes, filename: str) -> list[ParsedDocument]:
    """解析 TXT 文件"""
    try:
        text = file_bytes.decode("utf-8").strip()
    except UnicodeDecodeError as e:
        logger.error(f"TXT 文件编码错误: {filename}, {e}")
        raise ValueError(f"TXT 文件编码错误，仅支持 UTF-8: {filename}")

    if not text:
        logger.warning(f"TXT 文件内容为空: {filename}")
        return []

    return [ParsedDocument(content=text, doc_type="txt", source=filename)]


def _parse_markdown(file_bytes: bytes, filename: str) -> list[ParsedDocument]:
    """解析 Markdown 文件"""
    try:
        text = file_bytes.decode("utf-8").strip()
    except UnicodeDecodeError as e:
        logger.error(f"Markdown 文件编码错误: {filename}, {e}")
        raise ValueError(f"Markdown 文件编码错误，仅支持 UTF-8: {filename}")

    if not text:
        logger.warning(f"Markdown 文件内容为空: {filename}")
        return []

    return [ParsedDocument(content=text, doc_type="md", source=filename)]


def _parse_pdf(file_bytes: bytes, filename: str) -> list[ParsedDocument]:
    """解析 PDF 文件 - 每页独立文档"""
    from pypdf import PdfReader

    # 写入临时文件供 PdfReader 解析
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        reader = PdfReader(tmp_path)
        documents = []

        for page_num, page in enumerate(reader.pages, start=1):
            text = page.extract_text()
            if text and text.strip():
                documents.append(
                    ParsedDocument(
                        content=text.strip(),
                        doc_type="pdf",
                        source=filename,
                        page_number=page_num,
                        extra_metadata={"total_pages": len(reader.pages)},
                    )
                )
            else:
                logger.warning(f"PDF 页面 {page_num} 内容为空: {filename}")

        if not documents:
            logger.warning(f"PDF 文件无有效内容: {filename}")

        return documents
    except Exception as e:
        logger.error(f"PDF 解析失败: {filename}, {e}")
        raise ValueError(f"PDF 解析失败: {filename}, 错误: {str(e)}")
    finally:
        Path(tmp_path).unlink(missing_ok=True)


def _parse_docx(file_bytes: bytes, filename: str) -> list[ParsedDocument]:
    """解析 Word 文件"""
    from docx import Document

    # 写入临时文件供 Document 解析
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        doc = Document(tmp_path)
        # 提取所有段落文本
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        text = "\n\n".join(paragraphs)

        if not text:
            logger.warning(f"Word 文件内容为空: {filename}")
            return []

        return [ParsedDocument(content=text, doc_type="docx", source=filename)]
    except Exception as e:
        logger.error(f"Word 解析失败: {filename}, {e}")
        raise ValueError(f"Word 解析失败: {filename}, 错误: {str(e)}")
    finally:
        Path(tmp_path).unlink(missing_ok=True)


def _parse_xlsx(file_bytes: bytes, filename: str) -> list[ParsedDocument]:
    """解析 Excel 文件 - 每个 Sheet 独立文档"""
    from openpyxl import load_workbook

    # 写入临时文件供 load_workbook 解析
    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        wb = load_workbook(tmp_path, data_only=True)
        documents = []

        for sheet_idx, sheet_name in enumerate(wb.sheetnames, start=0):
            sheet = wb[sheet_name]
            # 提取所有行的文本
            rows_text = []
            for row in sheet.iter_rows(values_only=True):
                # 过滤空行，将单元格值转为字符串
                row_values = [str(cell) if cell is not None else "" for cell in row]
                row_text = " | ".join(row_values)
                if row_text.strip():
                    rows_text.append(row_text)

            text = "\n".join(rows_text)
            if text.strip():
                documents.append(
                    ParsedDocument(
                        content=text.strip(),
                        doc_type="xlsx",
                        source=filename,
                        sheet_name=sheet_name,
                        sheet_index=sheet_idx,
                        extra_metadata={"total_sheets": len(wb.sheetnames)},
                    )
                )
            else:
                logger.warning(f"Excel Sheet '{sheet_name}' 内容为空: {filename}")

        if not documents:
            logger.warning(f"Excel 文件无有效内容: {filename}")

        return documents
    except Exception as e:
        logger.error(f"Excel 解析失败: {filename}, {e}")
        raise ValueError(f"Excel 解析失败: {filename}, 错误: {str(e)}")
    finally:
        Path(tmp_path).unlink(missing_ok=True)


def _parse_pptx(file_bytes: bytes, filename: str) -> list[ParsedDocument]:
    """解析 PPT 文件 - 每页独立文档"""
    from pptx import Presentation

    # 写入临时文件供 Presentation 解析
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        prs = Presentation(tmp_path)
        documents = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            # 提取幻灯片中的所有文本
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())

            text = "\n\n".join(texts)
            if text.strip():
                documents.append(
                    ParsedDocument(
                        content=text.strip(),
                        doc_type="pptx",
                        source=filename,
                        page_number=slide_num,
                        extra_metadata={"total_slides": len(prs.slides)},
                    )
                )
            else:
                logger.warning(f"PPT 页面 {slide_num} 内容为空: {filename}")

        if not documents:
            logger.warning(f"PPT 文件无有效内容: {filename}")

        return documents
    except Exception as e:
        logger.error(f"PPT 解析失败: {filename}, {e}")
        raise ValueError(f"PPT 解析失败: {filename}, 错误: {str(e)}")
    finally:
        Path(tmp_path).unlink(missing_ok=True)


def validate_file(file_bytes: bytes, filename: str) -> tuple[bool, str]:
    """
    校验文件是否合法

    Args:
        file_bytes: 文件二进制内容
        filename: 文件名

    Returns:
        (是否合法, 错误信息或空字符串)
    """
    import filetype

    # 检查文件大小
    if len(file_bytes) > config.MAX_FILE_SIZE_BYTES:
        return False, f"文件大小超过限制 ({config.MAX_FILE_SIZE_MB}MB)"

    # 检查文件扩展名
    ext = Path(filename).suffix.lower()
    if ext not in config.EXT_TO_DOC_TYPE:
        return False, f"不支持的文件类型: {ext}"

    # 检查 MIME 类型（使用 filetype 库）
    # 注意：txt 和 md 文件可能无法被 filetype 正确识别，需要特殊处理
    if ext in (".txt", ".md"):
        # txt/md 文件不进行魔数检查，仅检查扩展名
        return True, ""

    # 其他文件类型使用 filetype 检查
    kind = filetype.guess(file_bytes)
    if kind is None:
        # filetype 无法识别，可能是文件损坏或格式不标准
        logger.warning(f"无法识别文件类型: {filename}, 扩展名: {ext}")
        # 不强制拒绝，允许继续尝试解析
        return True, ""

    # 检查 MIME 类型是否匹配
    expected_mimes = config.ALLOWED_FILE_TYPES.get(config.EXT_TO_DOC_TYPE[ext], [])
    if kind.mime not in expected_mimes:
        logger.warning(
            f"文件 MIME 类型不匹配: {filename}, "
            f"期望: {expected_mimes}, 实际: {kind.mime}"
        )
        # 不强制拒绝，允许继续解析（某些文件可能有非标准 MIME）

    return True, ""